from rest_framework import generics, permissions
from rest_framework.response import Response
from usr.serializers import *
from usr.models import *
from django.db.models import Q
from django.http import HttpResponse, JsonResponse
from rest_framework.views import APIView
from cmp.models import *
from ques.models import *
from ques.serializers import *
from cmp.serializers import * 
from com.models import *
from datetime import datetime,date
from datetime import time
from com.serializers import *
from BebrasBackend.constants import *
import requests 
from usr.password_encryption import encrypt,decrypt
import copy 
import six 
import time as t 
import os
from os import walk
import glob
from random import shuffle
from usr.password_encryption import encrypt,decrypt
from django.conf import settings
from django.http import HttpResponse
from django.http import FileResponse
from django.shortcuts import render
from rest_framework.views import APIView
from rest_framework.response import Response
import json
from pptx import Presentation
from rest_framework.decorators import parser_classes
from rest_framework.parsers import FileUploadParser,MultiPartParser
from PyPDF2 import PdfFileWriter, PdfFileReader
import shutil

class getCompetitionAPI(generics.GenericAPIView):
    permission_classes = [
      permissions.IsAuthenticated,
    ]
    def getTotalTime(studentenrolled):
        studentresponse=studentResponse.objects.filter(studentEnrollmentID=studentenrolled)
        print(studentresponse)
        totalTime=0.0
        for res in studentresponse:
          totalTime=totalTime+res.time
        totalTime=totalTime*60
        min, sec = divmod(totalTime, 60)
        hr, min = divmod(min, 60)
        totalTime = time(int(hr), int(min), int(sec))
        return totalTime
    def get(self, request):
      try:
        cmpNames=[]
        studentenrolled=studentEnrollment.objects.filter(userID=request.user.userID)
        print(studentenrolled)
        if not studentenrolled.exists():
          return Response("Not registered  for any  Competitions",status=404)
        for data in studentenrolled:
          endDate=data.competitionAgeID.competitionID.endDate
          startDate=data.competitionAgeID.competitionID.startDate
          if endDate.replace(tzinfo=None) > datetime.now()-timedelta(seconds=20) and datetime.now()-timedelta(seconds=20) >= startDate.replace(tzinfo=None):
            if data.score!=999:
              attempted=True
            else:
              attempted=False
            try:
              studentrespon=studentResponse.objects.filter(studentEnrollmentID=data)
              timetaken=getCompetitionAPI.getTotalTime(data)
              timeduration = datetime.combine(date.today(),data.competitionAgeID.competitionID.testDuration)-datetime.combine(date.today(),timetaken)
              timeduration= (datetime.min + timeduration).time()
                      
            except:
              timeduration=data.competitionAgeID.competitionID.testDuration
     
            cmpobj={"competitionname":data.competitionAgeID.competitionID.competitionName,"startDate":data.competitionAgeID.competitionID.startDate.date()
              ,"endDate":data.competitionAgeID.competitionID.endDate.date(),"testDuration":timeduration,"attempted":attempted}
            cmpNames.append(cmpobj)
        if len(cmpNames)==0:
          return Response("Sorry the competition you registered for has not started yet ",status=404)
        return  JsonResponse({"competitionnames":cmpNames}, safe=False)
      
      except Exception as e:
        return HttpResponse(e,status=404)

class getActiveCompetitionAPI(generics.GenericAPIView):
    permission_classes = [
      permissions.AllowAny,
    ]
  
    def get(self, request):
      try:
        cmpNames=[]
        competitions=competition.objects.filter(competitionType=main_challenge)
        for data in competitions:
          endDate=data.endDate
          startDate=data.startDate
          if endDate.year==datetime.now().year:
          # if endDate.replace(tzinfo=None) > datetime.now()-timedelta(seconds=20) and datetime.now()-timedelta(seconds=20) >= startDate.replace(tzinfo=None):
            cmpNames.append(data.competitionName)
        if len(cmpNames)==0:
          return Response("Sorry the competition you registered for has not started yet ",status=404)
        return  JsonResponse({"competitionnames":cmpNames}, safe=False)
      
      except Exception as e:
        return HttpResponse(e,status=404)

class getCmpQuestionAPI(generics.GenericAPIView):
    permission_classes = [
      permissions.IsAuthenticated,
    ]
    def getTotalTime(studentenrolled):
        studentresponse=studentResponse.objects.filter(studentEnrollmentID=studentenrolled)
        print(studentresponse)
        totalTime=0.0
        for res in studentresponse:
          totalTime=totalTime+res.time
        totalTime=totalTime*60
        min, sec = divmod(totalTime, 60)
        hr, min = divmod(min, 60)
        totalTime = time(int(hr), int(min), int(sec))
        return totalTime
    def post(self, request):
      try:
        testattempted=False
        response=[]
        print(request.data['competitionName'])
        compName=competition.objects.get(competitionName=request.data['competitionName'])
        studentenrolled=studentEnrollment.objects.filter(userID=request.user.userID).values_list('competitionAgeID', flat=True)
        studentenrolled=list(studentenrolled)
        if len(studentenrolled)==0:
          return Response("Not registered for any competitions ",status=404)
        try:
          cmpage = competitionAge.objects.get(competitionAgeID__in=studentenrolled,competitionID=compName.competitionID)
          print(cmpage)
        except:
          return Response("Either you are referring to a competition that student isn't enrolled for OR Same student has been  be registered for two different age groups  for the same competition",status=404)
        try: 
          studentEnrollmentlanguagecode=studentEnrollment.objects.get(competitionAgeID=cmpage.competitionAgeID,userID=request.user.userID)
        except:
          return Response("Same student has been  be registered twice for the same competition",status=404)
        
        cmpquestion=competitionQuestion.objects.filter(competitionAgeID=cmpage.competitionAgeID).values_list('questionID', flat=True)
        cmpquestion=list(cmpquestion)
        questiontranslation=questionTranslation.objects.filter(questionID__in=cmpquestion,languageCodeID=studentEnrollmentlanguagecode.languageCodeID)
        print("total number of questions",len(questiontranslation))
        questiondata={}
        for quest in questiontranslation:
          questiondata={}
          questioncaption=quest.translation['caption']
          questionbackground=quest.translation['background']
          questiondata['question_caption']=questioncaption
          questiondata['question_background']=questionbackground
          questiondata['question_domain']=quest.questionID.domainCodeID.codeName
          questiondata['identifier']=quest.Identifier
          question_skills=""
          cs_skills=quest.questionID.cs_skills.split(',')
          for skill in cs_skills:
            skill=int(skill)
            codeskills=code.objects.get(codeID=skill)
            print(codeskills)
            if question_skills=="":
              question_skills=codeskills.codeName
            else:
              question_skills=question_skills+" , "+codeskills.codeName
          questiondata['question_cs_skills']=question_skills
          cmpquest=competitionQuestion.objects.get(questionID=quest.questionID , competitionAgeID=cmpage.competitionAgeID)
          try:
             studentrespon=studentResponse.objects.get(studentEnrollmentID=studentEnrollmentlanguagecode,competitionQuestionID=cmpquest)
             if studentrespon.ansText==None and studentrespon.optionTranslationID==None:
               questiondata['answered']=False
             else:
               questiondata['answered']=True
             testattempted=True
          except:
             questiondata['answered']=False
             print(" no backup strategy") 
      
          cmpMarks=competition_MarkScheme.objects.get(competitionAgeID=cmpage.competitionAgeID,questionLevelCodeID=cmpquest.questionLevelCodeID)
          questiondata['question_marks']=cmpMarks.correctMarks
          correctoption=correctOption.objects.get(questionTranslationID=quest.questionTranslationID)
          if correctoption.ansText is None:
            if(quest.questionID.questionTypeCodeID.codeID==question_without_images):
              print("we came in question without images")
              questiondata['question_type']="question_without_images"
              optionlist=option.objects.filter(questionID=quest.questionID).values_list('optionID', flat=True)
              optionlist=list(optionlist)
              optiontranslation = optionTranslation.objects.filter(optionID__in=optionlist,languageCodeID=studentEnrollmentlanguagecode.languageCodeID.codeID)
              print("no of options",len(optiontranslation))
           
              optiondata=[]
              for optiond in optiontranslation:
                  optionjson={"optionTranslationID":optiond.optionTranslationID,"option":optiond.translationO['option']}
                  optiondata.append(optionjson)
              questiondata['options']=optiondata
              response.append(questiondata)
            elif(quest.questionID.questionTypeCodeID.codeID==question_with_images):
              print("we came in question with images")
              questiondata['question_type']="question_with_images"
              optionlist=option.objects.filter(questionID=quest.questionID).values_list('optionID', flat=True)
              optionlist=list(optionlist)
              optiontranslation=optionTranslation.objects.filter(optionID__in=optionlist,languageCodeID=studentEnrollmentlanguagecode.languageCodeID.codeID)
              print("no of options",len(optiontranslation))
              image=Image.objects.filter(ImageTypeCodeID=imageOpt)
              imagop=[]
              for im in image:
                for op in optiontranslation:
                  if(im.ObjectID==op.optionID.optionID): #image me question id ki question translation id?
                    imagop.append(im.uploadedFile)
              if(imagop):
                print("option has imagees")  
                questiondata['images_of_option']=imagop
              optiondata=[]
              for optiond in optiontranslation:
                  optionjson={"optionTranslationID":optiond.optionTranslationID,"option":optiond.translationO['option']}
                  optiondata.append(optionjson)
              questiondata['options']=optiondata
              response.append(questiondata)
          else:
            questiondata['question_type']="question_without_images_without_options"
            correctoption=correctOption.objects.get(questionTranslationID=quest.questionTranslationID)
            if (correctoption.ansText).isnumeric():
              questiondata['answertext_type']="number"
            else:
              questiondata['answertext_type']="text"
            response.append(questiondata)
        print("done")
        if testattempted:
          print("test attempted ")
          timeduration=getCmpQuestionAPI.getTotalTime(studentEnrollmentlanguagecode)
          timeduration = datetime.combine(date.today(),studentEnrollmentlanguagecode.competitionAgeID.competitionID.testDuration)-datetime.combine(date.today(),timeduration)
          timeduration= (datetime.min + timeduration).time()
         
        else:
          timeduration=studentEnrollmentlanguagecode.competitionAgeID.competitionID.testDuration
        studentEnrollmentlanguagecode.timeTaken=datetime.now()
        studentEnrollmentlanguagecode.save()
        shuffle(response)
        # print(response,type(response))
        print("Time updated")
        return  JsonResponse({"questions":response,"studentEnrollmentID":studentEnrollmentlanguagecode.studentEnrollmentID,"timeduration":timeduration}, safe=False)
      except Exception as e:
        return HttpResponse(status=404)

class getAlreadySavedResponse(generics.GenericAPIView):
    permission_classes = [
      permissions.IsAuthenticated,
    ]
    serializer_class = StudentResponseSerializer

    def post(self, request):
      try:
        selectedoption=""
        print(request.data)
        ques=request.data['identifier']
        quesTranslation=questionTranslation.objects.get(Identifier=ques)
        print("got question ",quesTranslation)
        student_enrolled=studentEnrollment.objects.get(studentEnrollmentID=request.data['studentEnrollmentID'])
        print(student_enrolled)
        cmpques=competitionQuestion.objects.get(questionID=quesTranslation.questionID.questionID,competitionAgeID=student_enrolled.competitionAgeID)
        print(cmpques.competitionQuestionID)
        opti=option.objects.filter(questionID=quesTranslation.questionID).values_list('optionID', flat=True)
        opti=list(opti)
        if len(opti)==0:
            request.data['optionTranslationID']=None
            request.data['ansText']=""
        else:
            request.data['ansText']=None
            request.data['optionTranslationID']=1
        del request.data['option']
        del request.data['identifier']
        print(request.data)   
        try:
          #student has already given response
          student_res=studentResponse.objects.get(competitionQuestionID=cmpques.competitionQuestionID,studentEnrollmentID=student_enrolled.studentEnrollmentID)
          if request.data['optionTranslationID']==None:
            print("Ans Text Type")
            if student_res.ansText==None:
              selectedoption=""
            else:
              selectedoption=student_res.ansText
          elif request.data['ansText']==None:
            if student_res.optionTranslationID==None:
              selectedoption=""
            else:
              selectedoption=student_res.optionTranslationID.translationO["option"]
          print(" Student Response Updated",selectedoption)
          return Response({"Option":selectedoption},status=200)
        except studentResponse.DoesNotExist:
          #student is giving response for the first time
          return Response({"Option":""},status=200)
   
      except Exception as e:
        return HttpResponse(e,status=404)

class studentResponseAPI(generics.GenericAPIView):
    permission_classes = [
      permissions.IsAuthenticated,
    ]
    serializer_class = StudentResponseSerializer

    def post(self, request):
      try:
        print(request.data)
        ques=request.data['identifier']
        quesTranslation=questionTranslation.objects.get(Identifier=ques)
        print("got question ",quesTranslation)
        student_enrolled=studentEnrollment.objects.get(studentEnrollmentID=request.data['studentEnrollmentID'])
        print(student_enrolled)
        # assuming only competitionageid and question id will be a unique pair
        cmpques=competitionQuestion.objects.get(questionID=quesTranslation.questionID.questionID,competitionAgeID=student_enrolled.competitionAgeID)
        print(cmpques.competitionQuestionID)
        opt=request.data['option']
        opti=option.objects.filter(questionID=quesTranslation.questionID).values_list('optionID', flat=True)
        opti=list(opti)
        if len(opti)==0:
            if request.data['option']=="":
              request.data['ansText']=None
            else:
              request.data['ansText']=request.data['option']
            request.data['optionTranslationID']=None
        else:
            if opt=="":
              request.data['optionTranslationID']=None
              request.data['ansText']=None
            else:
              optiontext['option']=opt
              opTranslation=optionTranslation.objects.get(optionTranslationID=opt)
              print(opTranslation)
              request.data['optionTranslationID']=opTranslation.optionTranslationID
              request.data['ansText']=None
        request.data['competitionQuestionID']=cmpques.competitionQuestionID
        startTime=student_enrolled.timeTaken #START TIME refers to student enrollment time taken
        endTime=datetime.now().time()        #END TIME refers to the current time
        time_difference = datetime.combine(date.today(),endTime)-datetime.combine(date.today(),startTime)
        time_diff=round(time_difference.seconds/60,4)
        request.data['time']=time_diff
        del request.data['option']
        del request.data['identifier']
        print(request.data)   
        try:
          #student has already given response
          student_res=studentResponse.objects.get(competitionQuestionID=cmpques.competitionQuestionID,studentEnrollmentID=student_enrolled.studentEnrollmentID)
          print(student_res,time_diff)
          student_res.time=round(student_res.time,4)+time_diff
          print(request.data['optionTranslationID'],request.data['ansText'])
          if request.data['optionTranslationID']==None and request.data['ansText']==None:
            student_res.optionTranslationID=None
            student_res.ansText=None
          elif  request.data['optionTranslationID']!=None and request.data['ansText']==None:
            student_res.optionTranslationID=opTranslation
          else:
            student_res.ansText=request.data['ansText']
            print(student_res.ansText)
          student_res.save()
          print(" Student Response Updated")
        except studentResponse.DoesNotExist:
          #student is giving response for the first time
          serializer = StudentResponseSerializer(data=request.data)
          if serializer.is_valid():
            studentres = serializer.save() 
            print(" Student Response Saved")
          else:
            errors=''
            for i in serializer.errors : 
              errors=errors+ i +' '+ str(serializer.errors[i])
              errors="Error : "+errors
            return Response(errors,status=404)
        student_enrolled.timeTaken=datetime.now()
        student_enrolled.save()
        print("Time updated")  
        return Response("Success",status=200)
   
      except Exception as e:
        return HttpResponse(e,status=404)

class validateOfflineUpload(generics.GenericAPIView):
    permission_classes = [
      permissions.AllowAny,
    ]
    serializer_class = StudentResponseSerializer
    def IsUserRegistered(data):
      print(data['userid'])
      try:
        if User.objects.filter(loginID=data['userid']).count() == 0:
            return False
        else:
            return True
      except Exception as e:
        print(e)
    def IsLoginCredsCorrect(data):
      print(data['userid'],data['paswd'])
      user=User.objects.get(loginID=data['userid'])
      pass_decrypt=decrypt(user.password)
      print(pass_decrypt)
      if(pass_decrypt==data['paswd'].strip()):
          return True
      else:
          return False
    def IsStudentEnrolled(data,compName):
      user=User.objects.get(loginID=data['userid'])
      studentenrollmentid=None
      student_enrolled=studentEnrollment.objects.filter(userID=user.userID)
      for studentenroll in student_enrolled:
        if studentenroll.competitionAgeID.competitionID.competitionName==compName:
          studentenrollmentid=studentenroll
      if studentenrollmentid==None:
          return False
      else:
          return True 
    def IsAgeGroupValid(data,compName):
      print(data['userid'],data['paswd'])
      grp= (data['group']).lower()
      studentenrollmentid=None
      for key, val in maps_groups.items(): 
          if val == grp: 
            ageGroup=key
      print(ageGroup)
      user=User.objects.get(loginID=data['userid'])
      student_enrolled=studentEnrollment.objects.filter(userID=user.userID)
      for studentenroll in student_enrolled:
        print(studentenroll.competitionAgeID.competitionID.competitionName,compName,studentenroll.competitionAgeID.AgeGroupClassID.AgeGroupID.AgeGroupName==ageGroup)
        if studentenroll.competitionAgeID.competitionID.competitionName==compName and studentenroll.competitionAgeID.AgeGroupClassID.AgeGroupID.AgeGroupName==ageGroup :
          studentenrollmentid=studentenroll
      if studentenrollmentid==None:
          return False
      else:
          return True
  
    def post(self, request):
      try:
        error=False
        print(type(request.data))
        
        compName=request.data['competitionName']
        responses=request.data['responses']

        for data in responses:
          responsestring=""
          print(data)
          newdata = dict((k.lower(), v) for k, v in data .items()) 
          print(newdata)
          result=validateOfflineUpload.IsUserRegistered(newdata)
          if(not(result)):
            responsestring=responsestring+"User Not Registered, UserID password do not match, User not enrolled, Agegroup Not Correct"
            data.update( {'Comments' : responsestring} )
            # data['Comments']=responsestring
            error=True
            continue
          print("Step 1 succesfull")
          result=validateOfflineUpload.IsLoginCredsCorrect(newdata)
          if(not(result)):
            responsestring=responsestring+"UserID password do not match"
          result=validateOfflineUpload.IsStudentEnrolled(newdata,compName)
          print("Step 2 succesfull")
          if(not(result)):
            responsestring=responsestring+", User not enrolled"
          result=validateOfflineUpload.IsAgeGroupValid(newdata,compName)
          print("Step 3 succesfull")
          if(not(result)):
            responsestring=responsestring+", Agegroup Not Correct"
          data.update( {'Comments' : responsestring} )
          if responsestring!="":
            error=True
        print(responses)
        if error:
            return JsonResponse(responses, safe=False,status=400)
        else:
            return Response("Success",status=200)


      except Exception as e:
        return HttpResponse(e,status=404)

class studentResponseFromExcelAPI(generics.GenericAPIView):
    permission_classes = [
      permissions.AllowAny,
    ]
    serializer_class = StudentResponseSerializer
    def calculateScore(studentenrollmentid):
        studentenrolled=studentEnrollment.objects.get(studentEnrollmentID=studentenrollmentid)
        studentenrolled.score=0
        studentenrolled.save()
        studentenrolled=studentEnrollment.objects.get(studentEnrollmentID=studentenrollmentid)
        studentresponses=studentResponse.objects.filter(studentEnrollmentID=studentenrolled)
        cmpquest=competitionQuestion.objects.filter(competitionAgeID=studentenrolled.competitionAgeID)
        print(len(studentresponses),len(cmpquest))
        if(len(studentresponses)!=len(cmpquest)):
          return "error"
        for data in cmpquest:
          questr=questionTranslation.objects.get(questionID=data.questionID,languageCodeID=studentenrolled.languageCodeID)
          correctOpt=correctOption.objects.get(questionTranslationID=questr.questionTranslationID)
          studentresponse=studentResponse.objects.get(studentEnrollmentID=studentenrolled,competitionQuestionID=data.competitionQuestionID)
          competition_Marks=competition_MarkScheme.objects.get(competitionAgeID=studentenrolled.competitionAgeID,questionLevelCodeID=studentresponse.competitionQuestionID.questionLevelCodeID)
          if correctOpt.optionTranslationID is not None:
            if studentresponse.optionTranslationID is not None:
              if correctOpt.optionTranslationID.optionTranslationID == studentresponse.optionTranslationID.optionTranslationID:
                print("correct")
                studentenrolled.score=studentenrolled.score+competition_Marks.correctMarks
                studentenrolled.save()
              elif correctOpt.optionTranslationID.optionTranslationID != studentresponse.optionTranslationID.optionTranslationID:
                print("Incorrect")
                studentenrolled.score=studentenrolled.score+competition_Marks.incorrectMarks
                studentenrolled.save()

          elif correctOpt.optionTranslationID is  None and correctOpt.ansText is not None:
            if studentresponse.ansText is not None:
              if correctOpt.ansText == studentresponse.ansText:
                print("correct")
                studentenrolled.score=studentenrolled.score+competition_Marks.correctMarks
                studentenrolled.save()
              elif correctOpt.ansText != studentresponse.ansText:
                print("Incorrect")
                studentenrolled.score=studentenrolled.score+competition_Marks.incorrectMarks
                studentenrolled.save()
    
    def post(self, request):
      try:
        error=False
        questionmappings=[]    
        answerlist=["A","B","C","D","E"]
        compName=request.data['competitionName']
        responses=request.data['responses']
        #validating data
        for data in responses:
          responsestring=""
          newdata = dict((k.lower(), v) for k, v in data .items()) 
          print(newdata)
          result=validateOfflineUpload.IsUserRegistered(newdata)
          if(not(result)):
            responsestring=responsestring+"User Not Registered, UserID password do not match, User not enrolled, Agegroup Not Correct"
            error=True
            continue
          result=validateOfflineUpload.IsLoginCredsCorrect(newdata)
          if(not(result)):
            responsestring=responsestring+"UserID password do not match"
          result=validateOfflineUpload.IsStudentEnrolled(newdata,compName)
          if(not(result)):
            responsestring=responsestring+", User not enrolled"
          result=validateOfflineUpload.IsAgeGroupValid(newdata,compName)
          if(not(result)):
            responsestring=responsestring+", Agegroup Not Correct"
          if responsestring!="":
            error=True
        if error:
            return Response("Please Validate again",status=400)

        ## find out what age groups the give competition has
        comp=competition.objects.get(competitionName=compName)
        competitionage=competitionAge.objects.filter(competitionID=comp).values_list("AgeGroupClassID",flat=True)
        Agegroupclasses=AgeGroupClass.objects.filter(AgeGroupClassID__in=list(competitionage)).values_list("AgeGroupID",flat=True)
        Agegroups=AgeGroup.objects.filter(AgeGroupID__in=list(Agegroupclasses)).values_list("AgeGroupName",flat=True)
        Agegroups=list(Agegroups)
        print(Agegroups)
        #find out competition questions associated to each age group
        competitionage=competitionAge.objects.filter(competitionID=comp)
        for data in competitionage:
          maps={}
          if data.AgeGroupClassID.AgeGroupID.AgeGroupName in Agegroups:
            key_group=maps_groups[data.AgeGroupClassID.AgeGroupID.AgeGroupName]
            cmp_ques=competitionQuestion.objects.filter(competitionAgeID=data).values_list('questionID',flat=True)
            splitarray=(data.AgeGroupClassID.AgeGroupID.AgeGroupName.split("-"))
            language=splitarray[1]
            language=code.objects.get(codeName=language)
            question_identifiers=questionTranslation.objects.filter(questionID__in=list(cmp_ques),languageCodeID=language).values_list("Identifier",flat=True)
            maps={key_group:list(question_identifiers)}
            questionmappings.append(maps)
            Agegroups.remove(data.AgeGroupClassID.AgeGroupID.AgeGroupName)            
        print(questionmappings)
        #process responses for each student
        for d in responses:
          studentenrollmentid=None
          newdata = dict((k.lower(), v) for k, v in d.items()) 
          print(newdata['userid'],newdata['paswd'])
          current_user=User.objects.get(loginID=newdata['userid'])
          current_studentenrolled=studentEnrollment.objects.filter(userID=current_user)
          grp= (newdata['group']).lower()
          print(d)
          del newdata['userid']
          del newdata['paswd']
          del newdata['group']
          if 'comments' in newdata.keys(): 
            del newdata['comments']
          else:
             print("comments absent")

          for key, val in maps_groups.items(): 
            if val == grp: 
              ageGroup=key
          for data in questionmappings:
            if data.get(grp)!=None:
              questions=data.get(grp)
              print("got ",questions)
          if questions==None:
            return Response("You haven't entered proper AgeGroup in the column Group")
          for studentenroll in current_studentenrolled:
            if studentenroll.competitionAgeID.competitionID.competitionName==compName and studentenroll.competitionAgeID.AgeGroupClassID.AgeGroupID.AgeGroupName==ageGroup :
              studentenrollmentid=studentenroll
          print(studentenrollmentid)
          if studentenrollmentid==None:
            responsestring="The user "+current_user.loginID+" has not been enrolled "
            return Response(responsestring,status=400)
          print(newdata)
          for key, value in newdata.items():
            studentresponsedata={}
            print( (key), value)
            index=int(key)-1
            if index>=len(questions):
              continue
            ques=questions[index]
            quesTranslation=questionTranslation.objects.get(Identifier=ques)
            print("got question ",quesTranslation)
            student_enrolled=studentEnrollment.objects.get(studentEnrollmentID=studentenrollmentid.studentEnrollmentID)
            cmpques=competitionQuestion.objects.get(questionID=quesTranslation.questionID.questionID,competitionAgeID=student_enrolled.competitionAgeID)
            opt=value
            opti=option.objects.filter(questionID=quesTranslation.questionID).values_list('optionID', flat=True)
            opti=list(opti)
            if len(opti)==0:
                if value=="":
                  studentresponsedata['ansText']=None
                else:
                  studentresponsedata['ansText']=value
                studentresponsedata['optionTranslationID']=None
            else:
                if opt=="":
                  studentresponsedata['optionTranslationID']=None
                  studentresponsedata['ansText']=None
                else:
                  if value.isnumeric():
                    studentresponsedata['optionTranslationID']=None
                  else:
                    position=answerlist.index(value.upper())
                    
                    optiontranslations=optionTranslation.objects.filter(optionID__in=opti,languageCodeID=student_enrolled.languageCodeID)
                    print("--------------------------------------",len(optiontranslations))
                    print(optiontranslations,value,position)
                    opTranslation=optionTranslation.objects.get(optionTranslationID=(optiontranslations[position]).optionTranslationID)
                    print(opTranslation)
                    studentresponsedata['optionTranslationID']=opTranslation.optionTranslationID
                  studentresponsedata['ansText']=None
            studentresponsedata['competitionQuestionID']=cmpques.competitionQuestionID
            # startTime=student_enrolled.timeTaken #START TIME refers to student enrollment time taken
            # endTime=datetime.now().time()        #END TIME refers to the current time
            # time_difference = datetime.combine(date.today(),endTime)-datetime.combine(date.today(),startTime)
            # time_diff=round(time_difference.seconds/60,4)
            time_diff=1.0
            studentresponsedata['time']=1.0
            print(studentresponsedata)   
            try:
              #student has already given response
              student_res=studentResponse.objects.get(competitionQuestionID=cmpques.competitionQuestionID,studentEnrollmentID=student_enrolled.studentEnrollmentID)
              student_res.time=round(student_res.time,4)+time_diff
              if studentresponsedata['optionTranslationID']==None and studentresponsedata['ansText']==None:
                student_res.optionTranslationID=None
                student_res.ansText=None
              elif  studentresponsedata['optionTranslationID']!=None and studentresponsedata['ansText']==None:
                student_res.optionTranslationID=opTranslation
              else:
                student_res.ansText=studentresponsedata['ansText']
                print(student_res.ansText)
              student_res.save()
              print(" Student Response Updated")
            except studentResponse.DoesNotExist:
              #student is giving response for the first time
              print("Saving now")
              studentresponsedata['studentEnrollmentID']=student_enrolled.studentEnrollmentID
              serializer = StudentResponseSerializer(data=studentresponsedata)
              if serializer.is_valid():
                studentres = serializer.save() 
                print(" Student Response Saved")
              else:
                errors=''
                for i in serializer.errors : 
                  errors=errors+ i +' '+ str(serializer.errors[i])
                  errors="Error : "+errors
                return Response(errors,status=404)
            student_enrolled.timeTaken=datetime.now()
            student_enrolled.save()
            print("Time updated")  
          statusdata=studentResponseFromExcelAPI.calculateScore(studentenrollmentid.studentEnrollmentID)
          if statusdata=="error":
            responsestring="You haven't completed the test for "+studentenrollmentid.userID.loginID
            return Response(responsestring,status=400)
        return Response("Success",status=200)
   
      except Exception as e:
        return HttpResponse(e,status=404)

class calcTotalScore(generics.GenericAPIView):
    permission_classes = [
      permissions.IsAuthenticated,
    ]
    serializer_class = CorrectOptionSerializer

    def post(self, request):
      try:
        print(request.user.userID)
        studentenrolled=studentEnrollment.objects.get(studentEnrollmentID=request.data['studentEnrollmentID'])
        studentenrolled.score=0
        studentenrolled.save()
        studentenrolled=studentEnrollment.objects.get(studentEnrollmentID=request.data['studentEnrollmentID'])
        cmpquest=competitionQuestion.objects.filter(competitionAgeID=studentenrolled.competitionAgeID)
        for data in cmpquest:
          questr=questionTranslation.objects.get(questionID=data.questionID,languageCodeID=studentenrolled.languageCodeID)
          correctOpt=correctOption.objects.get(questionTranslationID=questr.questionTranslationID)
          studentresponse=studentResponse.objects.get(studentEnrollmentID=studentenrolled,competitionQuestionID=data.competitionQuestionID)
          competition_Marks=competition_MarkScheme.objects.get(competitionAgeID=studentenrolled.competitionAgeID,questionLevelCodeID=studentresponse.competitionQuestionID.questionLevelCodeID)
          if correctOpt.optionTranslationID is not None:
            if studentresponse.optionTranslationID is not None:
              if correctOpt.optionTranslationID.optionTranslationID == studentresponse.optionTranslationID.optionTranslationID:
                print("correct")
                studentenrolled.score=studentenrolled.score+competition_Marks.correctMarks
                studentenrolled.save()
                print("Score Saved")
              elif correctOpt.optionTranslationID.optionTranslationID != studentresponse.optionTranslationID.optionTranslationID:
                print("Incorrect")
                studentenrolled.score=studentenrolled.score+competition_Marks.incorrectMarks
                studentenrolled.save()
                print(" Score Saved")

    
      
          elif correctOpt.optionTranslationID is  None and correctOpt.ansText is not None:
            if studentresponse.ansText is not None:
              if correctOpt.ansText == studentresponse.ansText:
                print("correct")
                studentenrolled.score=studentenrolled.score+competition_Marks.correctMarks
                studentenrolled.save()
                print("Score Saved")
              elif correctOpt.ansText != studentresponse.ansText:
                print("Incorrect")
                studentenrolled.score=studentenrolled.score+competition_Marks.incorrectMarks
                studentenrolled.save()
                print(" Score Saved")
        studentresponse=studentResponse.objects.filter(studentEnrollmentID=studentenrolled)
        totalTime=0.0
        for res in studentresponse:
          totalTime=totalTime+res.time
        totalTime=totalTime*60
        min, sec = divmod(totalTime, 60)
        hr, min = divmod(min, 60)
        totalTime = time(int(hr), int(min), int(sec))
        studentenrolled.timeTaken=totalTime
        print("Time updated",totalTime)
        studentenrolled.save()
        return JsonResponse("success", safe=False)
      except Exception as e:
        return HttpResponse(e,status=404)     

class getCompetitionsNamesForResultAPI(generics.GenericAPIView):
    permission_classes = [
      permissions.IsAuthenticated,
    ]
    def get(self, request):
      try:
        print("Received request from ",request.user.userID)
        cmpNames=[]
        studentenrolled=studentEnrollment.objects.filter(userID=request.user.userID)
        if not studentenrolled.exists():
          return Response("User hasn't registered for any of the competitions",status=404)
        print(studentenrolled)
        for data in studentenrolled:
          if data.competitionAgeID.competitionID.endDate.replace(tzinfo=None) < datetime.now()-timedelta(seconds=20):
         
                if data.score!=999:
                    cmpNames.append(data.competitionAgeID.competitionID.competitionName)
                    print(cmpNames)
                else:
                  print("Registered but not completed  ",data.competitionAgeID.competitionID.competitionName)
        cmpNames.reverse()
        if len(cmpNames)==0:
          return Response("Sorry you haven't appeared for any competitions",status=404)
        return  JsonResponse({"competitionnames":cmpNames}, safe=False)
      
      except Exception as e:
        return HttpResponse(e,status=404)

class getCompetitionsNamesForAnalysisAPI(generics.GenericAPIView):
    permission_classes = [
      permissions.IsAuthenticated,
    ]
    def get(self, request):
      try:
        print("Received request from ",request.user.userID)
        cmpNames=[]
        role_of_user=UserRole.objects.get(userID=request.user.userID)
        location_of_user=UserRoleLocation.objects.get(userRoleID=role_of_user)
        schoolname=school.objects.get(schoolID=location_of_user.locationObjectID)
        mainChallenge=code.objects.get(codeID=main_challenge)
        comp=competition.objects.filter(competitionType=mainChallenge)
        for data in comp:
          if data.endDate.replace(tzinfo=None) < datetime.now()-timedelta(seconds=20):
            if data.endDate > schoolname.registered_On:
                cmpNames.append(data.competitionName)
                print(cmpNames)
            else:
                print("Not appeared for ",data.competitionName)
        # cmpNames.reverse()
        if len(cmpNames)==0:
          return Response("Sorry your school hasn't appeared for any competitions",status=404)
        return  JsonResponse({"competitionnames":cmpNames}, safe=False)
      
      except Exception as e:
        return HttpResponse(e,status=404)

class getCmpResultsAPI(generics.GenericAPIView):
    permission_classes = [
      permissions.IsAuthenticated,
    ]
    def post(self, request):
      try:
        totalscore=0
        Resultsresponse=[]
        print("Request received from ",request.user.userID)
        compName=competition.objects.get(competitionName=request.data['competitionName'])
        print(compName)
        studentenrolled=studentEnrollment.objects.filter(userID=request.user.userID).values_list('competitionAgeID', flat=True)
        if not studentenrolled.exists():
          return Response("User hasn't registered for any of the competitions",status=404)
        studentenrolled=list(studentenrolled)
        print(studentenrolled)
        
        cmpage = competitionAge.objects.get(competitionAgeID__in=studentenrolled,competitionID=compName.competitionID)
        print(cmpage)
        try:
          studentEnrollmentlanguagecode=studentEnrollment.objects.get(competitionAgeID=cmpage.competitionAgeID,userID=request.user.userID)
        except:
          return Response("User has registered more than once for same competition",status=404)
        if studentEnrollmentlanguagecode.score==999:
          print("User hasn't completed the test")
          return Response("User hasn't completed the test",status=404)
        domains=code.objects.filter(codeGroupID=domain)
        for domainnames in domains:
          resultjson={"domain":"","questions":[],"score":0,"marks":0}
          resultjson["domain"]=domainnames.codeName
          score=0
          marks=0
          response=[]
          cmpquestion=competitionQuestion.objects.filter(competitionAgeID=cmpage.competitionAgeID).values_list('questionID', flat=True)
          cmpquestion=list(cmpquestion)
          questionslist=question.objects.filter(questionID__in=cmpquestion,domainCodeID=domainnames).values_list('questionID', flat=True)
          questionslist=list(questionslist)
          if len(questionslist)==0:
            continue
          questiontranslation=questionTranslation.objects.filter(questionID__in=questionslist,languageCodeID=studentEnrollmentlanguagecode.languageCodeID)
          print("total number of questions",len(questiontranslation))
          questiondata={}
          for quest in questiontranslation:
            questiondata={}
            questiondata['question_caption']=quest.translation['caption']
            questiondata['question_background']=quest.translation['background']
            questiondata['question_explanation']=quest.translation['explanation']      
            questiondata['question_domain']=quest.questionID.domainCodeID.codeName
            questiondata['identifier']=quest.Identifier
            question_skills=""
            cs_skills=quest.questionID.cs_skills.split(',')
            print(cs_skills)
            for skill in cs_skills:
              skill=int(skill)
              codeskills=code.objects.get(codeID=skill)
              print(codeskills)
              if question_skills=="":
                question_skills=codeskills.codeName
              else:
                question_skills=question_skills+" , "+codeskills.codeName
            questiondata['question_cs_skills']=question_skills
            correctoption=correctOption.objects.get(questionTranslationID=quest.questionTranslationID)
            cmpquest=competitionQuestion.objects.get(questionID=quest.questionID,competitionAgeID=cmpage.competitionAgeID)
            print(cmpquest.competitionQuestionID,studentEnrollmentlanguagecode.studentEnrollmentID)
            try:
              studentres=studentResponse.objects.get(competitionQuestionID=cmpquest.competitionQuestionID,studentEnrollmentID=studentEnrollmentlanguagecode.studentEnrollmentID)
              print(studentres)
              if correctoption.ansText is None:
                if(quest.questionID.questionTypeCodeID.codeID==question_without_images):
                  print("we came in question without images")
                  questiondata['question_type']="question_without_images"
                  optionlist=option.objects.filter(questionID=quest.questionID).values_list('optionID', flat=True)
                  optionlist=list(optionlist)
                  optiontranslation = optionTranslation.objects.filter(optionID__in=optionlist,languageCodeID=studentEnrollmentlanguagecode.languageCodeID.codeID)
                  print("no of options",len(optiontranslation))
                  optiondata=[]
                  if studentres.optionTranslationID is None:
                    questiondata['selectedoption']="No option Selected"
                  else:
                    print("Got user's selected option")
                    questiondata['selectedoption']=studentres.optionTranslationID.translationO['option']
                  for optiond in optiontranslation:
                      optiondata.append(optiond.translationO['option'])
                  questiondata['options']=optiondata
                  questiondata['correctoption']=correctoption.optionTranslationID.translationO['option']
                  cmpmarks=competition_MarkScheme.objects.get(competitionAgeID=cmpage.competitionAgeID,questionLevelCodeID=cmpquest.questionLevelCodeID)
                  questiondata['marks']=cmpmarks.correctMarks
                  marks=marks+cmpmarks.correctMarks
                  totalscore=totalscore+cmpmarks.correctMarks
                  if(questiondata['selectedoption']==questiondata['correctoption']):
                    questiondata['Score']=cmpmarks.correctMarks
                  elif questiondata['selectedoption']=="No option Selected":
                    questiondata['Score']=0
                  else:
                    questiondata['Score']=cmpmarks.incorrectMarks
                  score=score+questiondata['Score']
                  response.append(questiondata)
                elif(quest.questionID.questionTypeCodeID.codeID==question_with_images):
                  print("we came in question with images")
                  questiondata['question_type']="question_with_images"
                  optionlist=option.objects.filter(questionID=quest.questionID).values_list('optionID', flat=True)
                  optionlist=list(optionlist)
                  optiontranslation=optionTranslation.objects.filter(optionID__in=optionlist,languageCodeID=studentEnrollmentlanguagecode.languageCodeID.codeID)
                  print("no of options",len(optiontranslation))
                  image=Image.objects.filter(ImageTypeCodeID=imageOpt)
                  imagop=[]
                  for im in image:
                    for op in optiontranslation:
                      if(im.ObjectID==op.optionID.optionID): #image me question id ki question translation id?
                        imagop.append(im.uploadedFile)
                  if(imagop):
                    print("option has imagees")  
                    questiondata['images_of_option']=imagop
                  if studentres.optionTranslationID is None:
                    questiondata['selectedoption']="No option Selected"
                  else:
                    print("Got user's selected option")
                    questiondata['selectedoption']=studentres.optionTranslationID.translationO['option']
                  optiondata=[]
                  for optiond in optiontranslation:
                      optiondata.append(optiond.translationO['option'])
                  questiondata['options']=optiondata
                  questiondata['correctoption']=correctoption.optionTranslationID.translationO['option']
                  cmpmarks=competition_MarkScheme.objects.get(competitionAgeID=cmpage.competitionAgeID,questionLevelCodeID=cmpquest.questionLevelCodeID)
                  questiondata['marks']=cmpmarks.correctMarks
                  marks=marks+cmpmarks.correctMarks
                  totalscore=totalscore+cmpmarks.correctMarks
                  if(questiondata['selectedoption']==questiondata['correctoption']):
                    questiondata['Score']=cmpmarks.correctMarks
                  elif questiondata['selectedoption']=="No option Selected":
                    questiondata['Score']=0
                  else:
                    questiondata['Score']=cmpmarks.incorrectMarks
                  score=score+questiondata['Score']
                  print(questiondata)
                  response.append(questiondata)
                
              else:
                questiondata['question_type']="question_without_images_without_options"
                questiondata['correctoption']=correctoption.ansText
                questiondata['selectedoption']=studentres.ansText
                cmpmarks=competition_MarkScheme.objects.get(competitionAgeID=cmpage.competitionAgeID,questionLevelCodeID=cmpquest.questionLevelCodeID)
                questiondata['marks']=cmpmarks.correctMarks
                marks=marks+cmpmarks.correctMarks
                totalscore=totalscore+cmpmarks.correctMarks
                if correctoption.ansText == studentres.ansText:
                  questiondata['Score']=cmpmarks.correctMarks
                elif  studentres.ansText == "" or studentres.ansText==None:
                  questiondata['Score']=0
                else :
                  questiondata['Score']=cmpmarks.incorrectMarks
                score=score+questiondata['Score']
                response.append(questiondata)
              
            except studentResponse.DoesNotExist:
              if correctoption.ansText is None:
                if(quest.questionID.questionTypeCodeID.codeID==question_without_images):
                  print("we came in question without images")
                  questiondata['question_type']="question_without_images"
                  optionlist=option.objects.filter(questionID=quest.questionID).values_list('optionID', flat=True)
                  optionlist=list(optionlist)
                  optiontranslation = optionTranslation.objects.filter(optionID__in=optionlist,languageCodeID=studentEnrollmentlanguagecode.languageCodeID.codeID)
                  print("no of options",len(optiontranslation))
                  optiondata=[]
                  questiondata['selectedoption']="No option Selected"
                  for optiond in optiontranslation:
                      optiondata.append(optiond.translationO['option'])
                  questiondata['options']=optiondata
                  questiondata['correctoption']=correctoption.optionTranslationID.translationO['option']
                  cmpmarks=competition_MarkScheme.objects.get(competitionAgeID=cmpage.competitionAgeID,questionLevelCodeID=cmpquest.questionLevelCodeID)
                  questiondata['marks']=cmpmarks.correctMarks
                  marks=marks+cmpmarks.correctMarks
                  totalscore=totalscore+cmpmarks.correctMarks
                  questiondata['Score']=0
                  response.append(questiondata)
                elif(quest.questionID.questionTypeCodeID.codeID==question_with_images):
                  print("we came in question with images")
                  questiondata['question_type']="question_with_images"
                  optionlist=option.objects.filter(questionID=quest.questionID).values_list('optionID', flat=True)
                  optionlist=list(optionlist)
                  optiontranslation=optionTranslation.objects.filter(optionID__in=optionlist,languageCodeID=studentEnrollmentlanguagecode.languageCodeID.codeID)
                  print("no of options",len(optiontranslation))
                  image=Image.objects.filter(ImageTypeCodeID=imageOpt)
                  imagop=[]
                  for im in image:
                    for op in optiontranslation:
                      if(im.ObjectID==op.optionID.optionID): #image me question id ki question translation id?
                        imagop.append(im.uploadedFile)
                  if(imagop):
                    print("option has imagees")  
                    questiondata['images_of_option']=imagop
                  
                  questiondata['selectedoption']="No option Selected"
                  
                  optiondata=[]
                  for optiond in optiontranslation:
                      optiondata.append(optiond.translationO['option'])
                  questiondata['options']=optiondata
                  questiondata['correctoption']=correctoption.optionTranslationID.translationO['option']
                  cmpmarks=competition_MarkScheme.objects.get(competitionAgeID=cmpage.competitionAgeID,questionLevelCodeID=cmpquest.questionLevelCodeID)
                  print(cmpmarks)
                  questiondata['marks']=cmpmarks.correctMarks
                  marks=marks+cmpmarks.correctMarks
                  totalscore=totalscore+cmpmarks.correctMarks
                  questiondata['Score']=0
                  print(questiondata)
                  response.append(questiondata)
                
              else:
                questiondata['question_type']="question_without_images_without_options"
                questiondata['correctoption']=correctoption.ansText
                questiondata['selectedoption']="No option Selected"
                cmpmarks=competition_MarkScheme.objects.get(competitionAgeID=cmpage.competitionAgeID,questionLevelCodeID=cmpquest.questionLevelCodeID)
                questiondata['marks']=cmpmarks.correctMarks
                marks=marks+cmpmarks.correctMarks
                totalscore=totalscore+cmpmarks.correctMarks
                questiondata['Score']=0
                print(questiondata)
                response.append(questiondata)
              
            except:
              return Response("One or more responses for the same question",status=404)
          resultjson["questions"]=response
          resultjson["score"]=score
          resultjson["marks"]=marks
          Resultsresponse.append(resultjson)
          print("done")
        totalscore=totalscore+studentEnrollmentlanguagecode.bonusMarks
        total=studentEnrollmentlanguagecode.score+studentEnrollmentlanguagecode.bonusMarks
        return  JsonResponse({"domain_wise_questions":Resultsresponse,"studentID":request.user.loginID,"studentName":request.user.username,"TotalScore":total,"totalMarks":totalscore}, safe=False)
      except Exception as e:
        return HttpResponse(e,status=404)

class LanguageAPI(APIView):

    permission_classes = [
      permissions.IsAuthenticated,
    ]
    def post(self, request):
      try:
        print(request.data)
        lang=[]
        comp=competition.objects.get(competitionName=request.data['competitionName'])
        compAge=competitionAge.objects.filter(competitionID=comp)
        print(list(compAge))
        for data in compAge:
          print(data.AgeGroupClassID.ClassID.classID)
          if data.AgeGroupClassID.ClassID.classID==request.data['class_id']:
            langarr=data.AgeGroupClassID.AgeGroupID.AgeGroupName.split('-')
            print(langarr[1])
            lang.append(langarr[1])
         
        if len(lang)==0:
          return Response("No languages to show for the selected competition",status=404)    
        return JsonResponse({"languages":lang}, safe=False)
      except Exception as e:
        return HttpResponse(e,status=404)

class practiceChallengeNames(APIView):

    permission_classes = [
      permissions.AllowAny,
    ]
    def get(self, request):
      try:
        comp=competition.objects.filter(competitionType=practice_challenge).order_by('-endDate')
        challengesnames=[]
        agegrp=[]
        print(len(comp))
        for data in comp:
          startDate=data.startDate.date()
          today=datetime.now().date()
          endDate=data.endDate.date() 
          if today > startDate and today < endDate:
            challenges=data.competitionID
            break
        cmpage=competitionAge.objects.filter(competitionID=challenges)
        print(len(cmpage))
        for data in cmpage:
          agegrp.append(data.AgeGroupClassID.AgeGroupID)
          challengesnames.append(data.AgeGroupClassID.AgeGroupID.AgeGroupName)
        
        challengesnames = list(dict.fromkeys(challengesnames))
        agegrp = list(dict.fromkeys(agegrp))
        for i in range(len(agegrp)):
          print(agegrp[i].AgeGroupName)
          agegrpclass=AgeGroupClass.objects.filter(AgeGroupID=agegrp[i]).values_list('ClassID', flat=True)
          caption=str(list(agegrpclass))
          caption=caption.replace('[', '')
          caption=caption.replace(']', '')
          index=challengesnames.index(agegrp[i].AgeGroupName)
          challengesnames[index]=challengesnames[index]+ " ( Class "+caption +")"
          print(challengesnames[i],challengesnames)
        if len(challengesnames)==0:
          return Response("No Practice Challenge to show ",status=404)    
        return JsonResponse({"PracticeChallenges":challengesnames}, safe=False)
      except Exception as e:
        return HttpResponse(e,status=404)

class PracChallengeQuestionAPI(generics.GenericAPIView):
    permission_classes = [
      permissions.AllowAny,
    ]
    def post(self, request):
      try:
        response=[]
        position=request.data['AgeGroupName'].index('(')
        request.data['AgeGroupName']=request.data['AgeGroupName'][:position]
        request.data['AgeGroupName']=request.data['AgeGroupName'].strip()
        print(request.data['AgeGroupName'])
        lang=request.data['AgeGroupName'].split('-')    
        language=code.objects.get(codeName=lang[1])
        comp=competition.objects.filter(competitionType=practice_challenge).order_by('-endDate')
        print(len(comp))
        for data in comp:
          startDate=data.startDate.date()
          today=datetime.now().date()
          endDate=data.endDate.date() 
          if today > startDate and today < endDate:
            challenges=data.competitionName
            break
        compName=competition.objects.filter(competitionName=challenges)
        print(compName[0],compName)
        agegrp=AgeGroup.objects.filter(AgeGroupName=request.data['AgeGroupName'])
        agegrpselected=None
        for ages in agegrp:
          if ages.created_on.year==datetime.now().year:
            agegrpselected=ages
        agegrpclasses=AgeGroupClass.objects.filter(AgeGroupID=agegrpselected).values_list('AgeGroupClassID', flat=True)
        agegrpclasses=list(agegrpclasses)
        newcmp=None
        cmpage = competitionAge.objects.filter(competitionID=compName[0].competitionID,AgeGroupClassID__in=agegrpclasses).values_list('competitionAgeID', flat=True)
        print((cmpage))
        if len(cmpage)==2:     
          cmpage1=competitionAge.objects.get(competitionAgeID=cmpage[0])
          cmpage2=competitionAge.objects.get(competitionAgeID=cmpage[1])
          if cmpage1.AgeGroupClassID.AgeGroupID.AgeGroupName==cmpage2.AgeGroupClassID.AgeGroupID.AgeGroupName:
            newcmp=cmpage1
            print(newcmp)
        else:
          return Response("Something went wrong",status=404)
        cmpage=newcmp
        cmpquestion=competitionQuestion.objects.filter(competitionAgeID=cmpage.competitionAgeID).values_list('questionID', flat=True)
        cmpquestion=list(cmpquestion)
        print(len(cmpquestion))
        questiontranslation=questionTranslation.objects.filter(questionID__in=cmpquestion,languageCodeID=language)
        print("total number of questions",len(questiontranslation))
        questiondata={}
        for quest in questiontranslation:
          questiondata={}
          questiondata['question_caption']=quest.translation['caption']
          questiondata['question_background']=quest.translation['background']
          questiondata['question_explanation']=quest.translation['explanation']
          questiondata['question_domain']=quest.questionID.domainCodeID.codeName
          questiondata['identifier']=quest.Identifier
          question_skills=""
          cs_skills=quest.questionID.cs_skills.split(',')
          for skill in cs_skills:
            skill=int(skill)
            codeskills=code.objects.get(codeID=skill)
            print(codeskills)
            if question_skills=="":
              question_skills=codeskills.codeName
            else:
              question_skills=question_skills+" , "+codeskills.codeName
          questiondata['question_cs_skills']=question_skills
          correctoption=correctOption.objects.get(questionTranslationID=quest.questionTranslationID)
          cmpquest=competitionQuestion.objects.get(questionID=quest.questionID,competitionAgeID=cmpage.competitionAgeID)
        
          if correctoption.ansText is None:
            if(quest.questionID.questionTypeCodeID.codeID==question_without_images):
              print("we came in question without images")
              questiondata['question_type']="question_without_images"
              optionlist=option.objects.filter(questionID=quest.questionID).values_list('optionID', flat=True)
              optionlist=list(optionlist)
              optiontranslation = optionTranslation.objects.filter(optionID__in=optionlist,languageCodeID=language.codeID)
              print("no of options",len(optiontranslation))
              optiondata=[]
         
              for optiond in optiontranslation:
                  optiondata.append(optiond.translationO['option'])
              questiondata['options']=optiondata
              questiondata['correctoption']=correctoption.optionTranslationID.translationO['option']
              print(quest.questionID,cmpage.AgeGroupClassID.AgeGroupID)
    
              cmpmarks=competition_MarkScheme.objects.get(competitionAgeID=cmpage.competitionAgeID,questionLevelCodeID=cmpquest.questionLevelCodeID)
              questiondata['Correct_marks']=cmpmarks.correctMarks
              questiondata['Incorrect_marks']=cmpmarks.incorrectMarks
              response.append(questiondata)
            elif(quest.questionID.questionTypeCodeID.codeID==question_with_images):
              print("we came in question with images")
              questiondata['question_type']="question_with_images"
              optionlist=option.objects.filter(questionID=quest.questionID).values_list('optionID', flat=True)
              optionlist=list(optionlist)
              optiontranslation=optionTranslation.objects.filter(optionID__in=optionlist,languageCodeID=language.codeID)
              print("no of options",len(optiontranslation))
              image=Image.objects.filter(ImageTypeCodeID=imageOpt)
              imagop=[]
              for im in image:
                for op in optiontranslation:
                  if(im.ObjectID==op.optionID.optionID): #image me question id ki question translation id?
                    imagop.append(im.uploadedFile)
              if(imagop):
                print("option has imagees")  
                questiondata['images_of_option']=imagop
            
              optiondata=[]
              for optiond in optiontranslation:
                  optiondata.append(optiond.translationO['option'])
              questiondata['options']=optiondata
              questiondata['correctoption']=correctoption.optionTranslationID.translationO['option']
              cmpmarks=competition_MarkScheme.objects.get(competitionAgeID=cmpage.competitionAgeID,questionLevelCodeID=cmpquest.questionLevelCodeID)
              print(cmpmarks)
              questiondata['Correct_marks']=cmpmarks.correctMarks
              questiondata['Incorrect_marks']=cmpmarks.incorrectMarks
             
          
              print(questiondata)
              response.append(questiondata)
             
          else:
            questiondata['question_type']="question_without_images_without_options"
            questiondata['correctoption']=correctoption.ansText
            if (correctoption.ansText).isnumeric():
              questiondata['answertext_type']="number"
            else:
              questiondata['answertext_type']="text"
            cmpmarks=competition_MarkScheme.objects.get(competitionAgeID=cmpage.competitionAgeID,questionLevelCodeID=cmpquest.questionLevelCodeID)
            print(cmpmarks)
            questiondata['Correct_marks']=cmpmarks.correctMarks
            questiondata['Incorrect_marks']=cmpmarks.incorrectMarks
            response.append(questiondata)
        print("done")
        return  JsonResponse({"questions":response}, safe=False)
      except Exception as e:
        return HttpResponse(e,status=404)

class GetAgeGroupToppers(APIView):
    def get(self,request,**kwargs):
        mainChallenge=code.objects.get(codeID=main_challenge)
        competitions=competition.objects.filter(competitionType=mainChallenge).order_by('-endDate')
        print(competitions)
        for comp in competitions:
          if comp.endDate.date() < datetime.now().date():
            cmpcurrent=comp
            print("got",cmpcurrent)
            break
        stud=[]
        ageGroupDict=[]
        responsedata=[]
        ageGroupList = []
        cmpage=competitionAge.objects.filter(competitionID=cmpcurrent)
        for compage in cmpage:
          ageGroupList.append(compage.AgeGroupClassID.AgeGroupID.AgeGroupName)
          try:
            studentenrolled=studentEnrollment.objects.filter(~Q(score = 999),competitionAgeID=compage).order_by('-score','timeTaken')
            if len(studentenrolled)!=0:
              stud.append(studentenrolled[0].studentEnrollmentID)
          except:
            return Response(status=400)
        ageGroupList = list(dict.fromkeys(ageGroupList))
        print(ageGroupList)
        timecurrent=time(int(1), int(1), int(30))
        for data in ageGroupList:
          studenttop={'AgeGroup':data,'score':0,'time':timecurrent}
          ageGroupDict.append(studenttop)
        print(ageGroupDict)
        for d in ageGroupDict:
          topperdata={}
          for student in stud:
            studentenroll=studentEnrollment.objects.get(studentEnrollmentID=student)
            resdata={}
            resdata['studentName']=studentenroll.userID.username
            resdata['schoolName']=studentenroll.schoolClassID.schoolID.schoolName
            resdata['gender']=studentenroll.userID.gender.codeName
            resdata['time']=studentenroll.timeTaken
            resdata['rank']=1
            resdata['score']=studentenroll.score+studentenroll.bonusMarks
            resdata['AgeGroup']=studentenroll.competitionAgeID.AgeGroupClassID.AgeGroupID.AgeGroupName
            studentscore=studentenroll.score+studentenroll.bonusMarks
            if d['AgeGroup']==studentenroll.competitionAgeID.AgeGroupClassID.AgeGroupID.AgeGroupName and d['score']<studentscore:
                d['score']=studentenroll.score
                topperdata=resdata
                d['time']=studentenroll.timeTaken
            elif d['AgeGroup']==studentenroll.competitionAgeID.AgeGroupClassID.AgeGroupID.AgeGroupName and d['score']==studentscore and d['time']>studentenroll.timeTaken:
                d['time']=studentenroll.timeTaken
                topperdata=resdata
            elif d['AgeGroup']==studentenroll.competitionAgeID.AgeGroupClassID.AgeGroupID.AgeGroupName and d['score']==studentscore and d['time']==studentenroll.timeTaken:
                topperdata=resdata
          if topperdata!={}:
            responsedata.append(topperdata)
        return Response({"toppers":responsedata,"competitionName":cmpcurrent.competitionName},status=200)

class GetSchoolToppers(APIView):
    def getTotalscore(competitionAgeid):
      totalscore=0
      compques=competitionQuestion.objects.filter(competitionAgeID=competitionAgeid)
      for ques in compques:
          competition_MarkS=competition_MarkScheme.objects.get(competitionAgeID=competitionAgeid,questionLevelCodeID=ques.questionLevelCodeID)
            
          totalscore=totalscore+competition_MarkS.correctMarks
      return totalscore
    def post(self,request,**kwargs):
        token=request.META.get('HTTP_AUTHORIZATION')
        c = CustomizePPT()
        t= GetLatestTemplate()
        g=GetStudentsAgeGroupWise()
        d=g.post(request,json=request.data, format='json',headers={'Authorization': token})
        template=t.get(request,type='School_Toppers')
        print('Latest',template.data)
        print((d.data))
        if len(d.data) == 0:
          return Response("no students from this class were enrolled for the competition",status=404)
        data=[]
    
        for i in range(0, len(d.data)):
            compage=competitionAge.objects.get(competitionAgeID=d.data[i]['competitionAgeID'])
            userdata=User.objects.get(userID=d.data[i]['userID'])
            schoolclassdata=schoolClass.objects.get(schoolClassID=d.data[i]['schoolClassID'])
            score=d.data[i]['score']+d.data[i]['bonusMarks']
            totalscore=GetParticipationCertificates.getTotalscore(d.data[i]['competitionAgeID'])+d.data[i]['bonusMarks']
            data.append({'year':compage.competitionID.startDate.year,
                         'Name':userdata.username,
                         'loginID':userdata.loginID,
                         'group':compage.AgeGroupClassID.AgeGroupID.AgeGroupName,
                         'score':str(score),
                         'total':totalscore,
                         'time': d.data[i]['timeTaken'],
                         'class':str(schoolclassdata.classNumber)})
       

     
        school=schoolclassdata.schoolID.schoolName+', '+schoolclassdata.schoolID.addressID.city
        type='schoolToppers'
        c.ppt(os.path.join(settings.MEDIA_ROOT) + '/ppt//'+template.data, data, school,type,duplicate=False)
        path_response=  school + '-Toppers-' + data[0][ 'group'] + '-' +str(data[0]['year'])
        return Response({"path":path_response},status=200)

class GetSchoolClassStudents(APIView):
    def post(self,request,**kwargs):
        print(request.user.userID)
        cmpname=request.data['cmpName']
        usrrole=UserRole.objects.get(userID=request.user.userID)
        usrroleLocation=UserRoleLocation.objects.get(userRoleID=usrrole.userRoleID)
        school_name = school.objects.get(schoolID=usrroleLocation.locationObjectID)
        comp= competition.objects.get(competitionName=cmpname)
        schclass=schoolClass.objects.get(schoolID=school_name.schoolID,classNumber=request.data['class_id'])
        cmpAge= competitionAge.objects.filter(competitionID=comp.competitionID )
        print(cmpAge)
        for data in cmpAge:
          print(data.AgeGroupClassID.ClassID.classID,request.data['class_id'])
          if data.AgeGroupClassID.ClassID.classID==request.data['class_id']:
            compage=data
            print(compage)
            break
        lists = studentEnrollment.objects.filter(~Q(score = 999),schoolClassID=schclass,competitionAgeID=compage.competitionAgeID)
        print(len(lists))
        if len(lists)==0:
          return Response("no students from this class were enrolled for the competition",status=404)
        serializer= studentEnrollmentViewSerializer(lists,many=True)
        return Response(serializer.data)

class GetParticipationCertificatesforStudentPortal(APIView):
    permission_classes = [
      permissions.IsAuthenticated,
    ]
    def post(self, request,**kwargs ):
        print(request.user.created_by,request.data['competitionName'])
        current_user=User.objects.get(loginID=request.user.created_by)
        c = CustomizePPT()
        usrrole=UserRole.objects.get(userID=current_user.userID)
        usrroleLocation=UserRoleLocation.objects.get(userRoleID=usrrole.userRoleID)
        school_name = school.objects.get(schoolID=usrroleLocation.locationObjectID)
        studentsenroll=studentEnrollment.objects.filter(userID=request.user.userID)
        for students in studentsenroll:
          if students.competitionAgeID.competitionID.competitionName==request.data['competitionName']:
            data_user=students
            break
        t=GetLatestTemplate()
        print("Here in certificates for Student Portal",data_user)
        data=[]
        print(data_user.competitionAgeID)
        totalscore=GetSchoolToppers.getTotalscore(data_user.competitionAgeID)+data_user.bonusMarks
        template=t.get(request,type='Participation')
        print('Latest',template.data)
      
        data.append({'year':data_user.competitionAgeID.competitionID.startDate.year,
                         'Name':request.user.username,
                         'loginID':request.user.loginID,
                         'group':data_user.competitionAgeID.AgeGroupClassID.AgeGroupID.AgeGroupName,
                         'score':str(data_user.score+data_user.bonusMarks),
                         'total':totalscore,
                         'class':str(data_user.competitionAgeID.AgeGroupClassID.ClassID.classID)})
           
       
        schooldata=school_name.schoolName+', '+school_name.addressID.city
        type='participation'
        path_response=schooldata+'-Class-'+data[0]['class']+'-'+data[0]['group']+'-'+str(data[0]['year'])
        c.ppt(os.path.join(settings.MEDIA_ROOT) + '/ppt//'+template.data, data, schooldata,type,duplicate=False)
        return Response({"path":path_response},status=200)

class GetParticipationCertificatesforSelectedStudents(APIView):
    permission_classes = [
      permissions.IsAuthenticated,
    ]
    def post(self, request,**kwargs ):
        print(request.user.userID)
        c = CustomizePPT()
        usrrole=UserRole.objects.get(userID=request.user.userID)
        usrroleLocation=UserRoleLocation.objects.get(userRoleID=usrrole.userRoleID)
        school_name = school.objects.get(schoolID=usrroleLocation.locationObjectID)
       
        t=GetLatestTemplate()
        data_user=request.data['users']
        print("Here in cert",data_user)
        data=[]
     
       
        template=t.get(request,type='Participation')
        print('Latest',template.data)
        for i in data_user:
            print(i['loginID'])
            comp= competition.objects.get(competitionName=i['competitionName'])
            data.append({'year':comp.startDate.year,
                         'Name':i['username'],
                         'loginID':i['loginID'],
                         'group':i['ageGroup'],
                         'score':str(i['score']),
                         'total':i['Totalscore'],
                         'class':str(i['class'])})
           
       
        schooldata=school_name.schoolName+', '+school_name.addressID.city
        type='participation'
        path_response=schooldata+'-Class-'+data[0]['class']+'-'+data[0]['group']+'-'+str(data[0]['year'])

        c.ppt(os.path.join(settings.MEDIA_ROOT) + '/ppt//'+template.data, data, schooldata,type,duplicate=True)
        return Response({"path":path_response},status=200)

class GetParticipationCertificates(APIView):

    def getTotalscore(competitionAgeid):
      print(competitionAgeid)
      totalscore=0
      compques=competitionQuestion.objects.filter(competitionAgeID=competitionAgeid)
      for ques in compques:
          competition_MarkS=competition_MarkScheme.objects.get(competitionAgeID=competitionAgeid,questionLevelCodeID=ques.questionLevelCodeID)   
          totalscore=totalscore+competition_MarkS.correctMarks
      return totalscore
    def post(self, request,**kwargs ):
        token=request.META.get('HTTP_AUTHORIZATION')
        c = CustomizePPT()
        g=GetSchoolClassStudents()
        t=GetLatestTemplate()
        d=g.post(request,json=request.data, format='json',headers={'Authorization': token})
        print("Here in cert",d.data)
        if d.data == "no students from this class were enrolled for the competition":
          return Response("no students from this class were enrolled for the competition",status=404)
        data=[]
        template=t.get(request,type='Participation')
        print('Latest',template.data)
        for i in range(0,len(d.data)):
            compage=competitionAge.objects.get(competitionAgeID=d.data[i]['competitionAgeID'])
            userdata=User.objects.get(userID=d.data[i]['userID'])
            schoolclassdata=schoolClass.objects.get(schoolClassID=d.data[i]['schoolClassID'])
            score=d.data[i]['score']+d.data[i]['bonusMarks']
            totalscore=GetParticipationCertificates.getTotalscore(d.data[i]['competitionAgeID'])+d.data[i]['bonusMarks']
            data.append({'year':compage.competitionID.startDate.year,
                         'Name':userdata.username,
                          'loginID':userdata.loginID,
                         'group':compage.AgeGroupClassID.AgeGroupID.AgeGroupName,
                         'score':str(score),
                         'total':totalscore,
                         'class':str(schoolclassdata.classNumber)})
         
        print((schoolclassdata.schoolID.schoolName))
        school=schoolclassdata.schoolID.schoolName+', '+schoolclassdata.schoolID.addressID.city
        type='participation'
        path_response=school+'-Class-'+data[0]['class']+'-'+data[0]['group']+'-'+str(data[0]['year'])

        c.ppt(os.path.join(settings.MEDIA_ROOT) + '/ppt//'+template.data, data, school,type,duplicate=True)
        return Response({"path":path_response},status=200)

class GetStudentsAgeGroupWise(APIView):
    def post(self,request,**kwargs):
        print(request.user.userID,request.data['class_id'])
        cmpname=request.data['cmpName']
        comp= competition.objects.get(competitionName=cmpname)
        usrrole=UserRole.objects.get(userID=request.user.userID)
        usrroleLocation=UserRoleLocation.objects.get(userRoleID=usrrole.userRoleID)
        school_name = school.objects.get(schoolID=usrroleLocation.locationObjectID)
        schclass=schoolClass.objects.get(schoolID=school_name.schoolID,classNumber=request.data['class_id'])
        cmpAge= competitionAge.objects.filter(competitionID=comp.competitionID )
        for data in cmpAge:
          print(data.AgeGroupClassID.ClassID.classID,request.data['class_id'])
          if data.AgeGroupClassID.ClassID.classID==request.data['class_id']:
            compage=data
            print(compage)
            break
        agegrps=AgeGroup.objects.filter(AgeGroupName=compage.AgeGroupClassID.AgeGroupID.AgeGroupName)
        agegroupid=0
        for i in range(0,len(agegrps)):
            if str(comp.startDate.year) in str(agegrps[i].created_on):
                agegroupid=agegrps[i].AgeGroupID
                break
        agegrpclasses=AgeGroupClass.objects.filter(AgeGroupID=agegroupid)
        agegroupclasses=[]
        for i in range(0,len(agegrpclasses)):
            agegroupclasses.append(agegrpclasses[i].AgeGroupClassID)
        cmpage=competitionAge.objects.filter(AgeGroupClassID__in= agegroupclasses)
        print(cmpage)
        cmpages=[]
        for i in range(0, len(cmpage)):
            cmpages.append(cmpage[i].competitionAgeID)

        schoolclasses = schoolClass.objects.filter(schoolID=school_name.schoolID)
        schclasses = []
        for i in range(0, len(schoolclasses)):
            schclasses.append(schoolclasses[i].schoolClassID)
        lists = studentEnrollment.objects.filter(~Q(score = 999),competitionAgeID__in =cmpages,schoolClassID__in=schclasses).order_by('-score','timeTaken')
        serializers=studentEnrollmentViewSerializer(lists,many=True)
        return Response(serializers.data)
        
@parser_classes((MultiPartParser,))
class CustomizePPT(APIView):
    def ppt(self,f,data,school,type,duplicate):
        prs = Presentation(f)
        NamesOfParticipants = []
        NOP = []
        text_runs = []
        def duplicate_slide(pres, index):
            template = pres.slides[index]
            try:
                blank_slide_layout = pres.slide_layouts[6]
            except:
                blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts) - 1]
            copied_slide = pres.slides.add_slide(blank_slide_layout)
            for shp in template.shapes:
                el = shp.element
                newel = copy.deepcopy(el)
                copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
            for _, value in six.iteritems(template.part.rels):
                if "notesSlide" not in value.reltype:
                    copied_slide.part.rels.add_relationship(value.reltype,
                                                            value._target,
                                                            value.rId)
            return copied_slide
        if duplicate==True:
            for i in range(0,len(data)-1):
                duplicate_slide(prs, 0)
        i = -1
        print("Count",len(prs.slides))
        for slide in prs.slides:
            if i < (len(data) - 1):
                i = i + 1
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                              if (shape.text.lower().find('name')) != -1:
                                  if(run.text.lower().find('name')!=-1):
                                    start=run.text.lower().find('name')
                                    end=run.text.lower().find('name')+len('name')
                                    nameofPDF=data[i]['Name']+'-'+data[i]['loginID']+'-'+str(data[i]['year'])
                                    NamesOfParticipants.append(nameofPDF)
                                    cur_text = run.text
                                    new_text = cur_text.replace(run.text[start:end], str(data[i]['Name']))
                                    run.text = new_text
                              if (shape.text.lower().find('year')) != -1:
                                  if (run.text.lower().find('year') != -1):
                                    start = run.text.lower().find('year')
                                    end = run.text.lower().find('year') + len('year')
                                    cur_text = run.text
                                    new_text = cur_text.replace(run.text[start:end], str(data[i]['year']))
                                    run.text = new_text
                              if (shape.text.lower().find('class')) != -1:
                                  if (run.text.lower().find('class') != -1):
                                    start = run.text.lower().find('class')
                                    end = run.text.lower().find('class') + len('class')
                                    cur_text = run.text
                                    new_text = cur_text.replace(run.text[start:end], str(data[i]['class']))
                                    run.text = new_text
                              if (shape.text.lower().find('group')) != -1:
                                  if (run.text.lower().find('group') != -1):
                                    start = run.text.lower().find('group')
                                    end = run.text.lower().find('group') + len('group')
                                    cur_text = run.text
                                    new_text = cur_text.replace(run.text[start:end], str(data[i]['group']))
                                    run.text = new_text
                              if (shape.text.lower().find('score')) != -1:
                                  if (run.text.lower().find('score') != -1):
                                    start = run.text.lower().find('score')
                                    end = run.text.lower().find('score') + len('score')
                                    cur_text = run.text
                                    new_text = cur_text.replace(run.text[start:end], str(data[i]['score']))
                                    run.text = new_text
                              if (shape.text.lower().find('total')) != -1:
                                  if (run.text.lower().find('total') != -1):
                                    start = run.text.lower().find('total')
                                    end = run.text.lower().find('total') + len('total')
                                    cur_text = run.text
                                    new_text = cur_text.replace(run.text[start:end], str(data[i]['total']))
                                    run.text = new_text  
                              if (shape.text.find('[')) != -1:
                                    cur_text = run.text
                                    new_text = cur_text.replace(str('['), str(''))
                                    run.text = new_text
                              if (shape.text.find(']')) != -1:
                                    cur_text = run.text
                                    new_text = cur_text.replace(str(']'), str(''))
                                    run.text = new_text
            else:
                rId = prs.slides._sldIdLst[-1].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[-1]
        if(type=='participation'):
          # NOP = NamesOfParticipants[::2]
          NOP =  list(dict.fromkeys(NamesOfParticipants))
        elif(type=='schoolToppers'):
          NOP = NamesOfParticipants
        path=''
        ldir = os.path.join(settings.MEDIA_ROOT) + '/output//'
        if(len(ldir)!=0):
          shutil.rmtree(ldir)
          os.makedirs(os.path.join(settings.MEDIA_ROOT) +'/output//')
        if(type=='participation'):
            path=os.path.join(settings.MEDIA_ROOT) + '/output//'+school+'-Class-'+data[0]['class']+'-'+data[0]['group']+'-'+str(data[0]['year'])+'.pptx'
        elif(type=='schoolToppers'):
            path=os.path.join(settings.MEDIA_ROOT) + '/output//' + school + '-Toppers-' + data[0][ 'group'] + '-' +str(data[0]['year'])+'.pptx'
        elif (type == 'nationalToppers'):
            path = os.path.join(settings.MEDIA_ROOT) + '/output//' +  'National Toppers-' + data[0]['group'] + '-' + data[0]['year'] + '.pptx'
        prs.save(path)
        # pythoncom.CoInitialize()
        def zipdir(self,source):
            print("Entered zipdir")
            pdfName=''
            if(type=='participation'):
                pdfName=school+'-Class-'+data[0]['class']+'-'+data[0]['group']+'-'+str(data[0]['year'])
            elif(type=='schoolToppers'):
                pdfName=school + '-Toppers-' + data[0][ 'group'] + '-' +str(data[0]['year'])
            basedir = pdfName
            baseName = os.path.join(settings.MEDIA_ROOT) + '/ZipFolder/BebrasCertificate'
            shutil.make_archive(base_dir=basedir,root_dir=source,format='zip',base_name=baseName)
            shutil.rmtree(os.path.join(settings.MEDIA_ROOT) + '/ZipFolder//'+pdfName)
            print("Done!")
        def splitPDF(self,path1):
            pdfName=''
            if(type=='participation'):
                pdfName=school+'-Class-'+data[0]['class']+'-'+data[0]['group']+'-'+str(data[0]['year'])
                CertiFolder = os.path.join(settings.MEDIA_ROOT)+'/ZipFolder//'+pdfName
                os.makedirs(CertiFolder)
            elif(type=='schoolToppers'):
                pdfName=school + '-Toppers-' + data[0][ 'group'] + '-' +str(data[0]['year'])
                CertiFolder = os.path.join(settings.MEDIA_ROOT)+'/ZipFolder//'+pdfName
                os.makedirs(CertiFolder)
            pdfFileObj = open(path1, "rb")
            inputpdf = PdfFileReader(pdfFileObj)
            i=0
            print(inputpdf.numPages)
            if(inputpdf.numPages==1):
                output = PdfFileWriter()
                output.addPage(inputpdf.getPage(i))
                outputFileName = os.path.join((settings.MEDIA_ROOT)+'/ZipFolder//'+pdfName,NOP[0]+".pdf" )
                with open(outputFileName, "wb") as outputStream:
                  output.write(outputStream)
            else:
                for i in range(inputpdf.numPages):
                    output = PdfFileWriter()
                    output.addPage(inputpdf.getPage(i))
                    outputFileName = os.path.join((settings.MEDIA_ROOT)+'/ZipFolder//'+pdfName,NOP[i]+".pdf")
                    with open(outputFileName, "wb") as outputStream:
                      output.write(outputStream)
            pdfFileObj.close()
   
        files = glob.glob(path)  # <--- ONLY CHANGE
        for filename in files:
                command = "unoconv -f pdf '" + filename+"'"
                os.system(command)
                os.remove(path)
        ldir =(os.path.join(settings.MEDIA_ROOT) + '/ZipFolder//')
        if(len(ldir)!=0):
          shutil.rmtree(ldir)
          os.makedirs(os.path.join(settings.MEDIA_ROOT) +'/ZipFolder//')
        path1=''
        if(type=='participation'):
            path1=os.path.join(settings.MEDIA_ROOT) + '/output//'+school+'-Class-'+data[0]['class']+'-'+data[0]['group']+'-'+str(data[0]['year'])+'.pdf'
        elif(type=='schoolToppers'):
            path1=os.path.join(settings.MEDIA_ROOT) + '/output//' + school + '-Toppers-' + data[0][ 'group'] + '-' +str(data[0]['year'])+'.pdf'
        print(path1,path)
        splitPDF(self,path1)
        source= os.path.join(settings.MEDIA_ROOT) + '/ZipFolder//'
        zipdir(self,source)
        
class deleteFiles(APIView):
    c = CustomizePPT()
    ppt = c.ppt
    def post(self,request,**kwargs):
      
        type=request.data['CertificateType']
        if(type=='participation'):
            os.remove(os.path.join(settings.MEDIA_ROOT) + '/output//'+request.data["path"]+ '.pdf')
        if(type=='schoolToppers'):
            os.remove(os.path.join(settings.MEDIA_ROOT) +  '/output//'+request.data["path"]+ '.pdf')
       
        os.remove(os.path.join(settings.MEDIA_ROOT) + '/ZipFolder/BebrasCertificate'+'.zip')
        return Response(status=200)

class GetLatestTemplate(APIView):               #Get Latest Template API
    permission_classes = (permissions.IsAuthenticated,)

    def get(self,request,**kwargs):
        type=kwargs['type']
        f = []
        for (dirpath, dirnames, filenames) in walk(os.path.join(settings.MEDIA_ROOT) + ("/ppt/")):
            for file in filenames:
                if (file.find(type) != -1):
                    f.append(file)
            break
        datetimes=[]
        for i in range(0,len(f)):
            arr=f[i].split('_')
            datetimes.append(arr[len(arr)-2]+'_'+arr[len(arr)-1].split('.')[0])
            sortedArray = sorted(
            datetimes,
            key=lambda x: datetime.strptime(x, '%Y-%m-%d_%H-%M-%S'), reverse=True
            )
        for i in range(0,len(f)):
            if str(sortedArray[0]) in f[i]:
                name=f[i]
                break
        return Response(name)
